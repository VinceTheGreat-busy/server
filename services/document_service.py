import io
import os

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


async def extract_document(filename: str, file_bytes: bytes):

    extension = os.path.splitext(filename)[1].lower()

    if extension == ".pdf":
        return extract_pdf(file_bytes)

    elif extension == ".docx":
        return extract_docx(file_bytes)

    elif extension == ".pptx":
        return extract_pptx(file_bytes)

    else:
        raise ValueError("Unsupported file type.")


def extract_pdf(file_bytes: bytes):

    pdf = PdfReader(io.BytesIO(file_bytes))

    text = []

    for page in pdf.pages:
        page_text = page.extract_text()

        if page_text:
            text.append(page_text)

    return "\n".join(text)


def extract_docx(file_bytes: bytes):

    document = Document(io.BytesIO(file_bytes))

    text = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            text.append(paragraph.text)

    return "\n".join(text)


def extract_pptx(file_bytes: bytes):

    presentation = Presentation(io.BytesIO(file_bytes))

    text = []

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)
